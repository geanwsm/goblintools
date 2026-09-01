import csv
import logging
import os
import re
import shutil
import subprocess
import unicodedata
from typing import Callable, Dict, List, Optional, Set, Tuple
from pathlib import Path
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from dbfread import DBF
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from pypdf.generic import IndirectObject

from goblintools.pypdf_workarounds import apply_pypdf_extraction_workarounds

apply_pypdf_extraction_workarounds()
import openpyxl
import xlrd
from odf import text, teletype
from odf.opendocument import load
from odf.text import P
import docx
from goblintools.ocr_parser import OCRProcessor
from goblintools.config import GoblinConfig, OCRConfig
from goblintools.log_policy import _set_suppress_warnings, log_warning
from goblintools.file_handling import FileValidator
from goblintools.ptbr_words import dict_hit_rate, is_probably_portuguese
from goblintools.extraction_report import (
    CLEAN,
    CORRUPT_UNRECOVERABLE,
    RECOVERED,
    ExtractionReport,
    PageExtraction,
)
from goblintools.table_extractor import (
    extract_pdf_tables,
    format_tables_for_page,
    tables_by_page,
)

logger = logging.getLogger(__name__)


def _has_meaningful_text(text: str) -> bool:
    """True if *text* has at least one letter, number, punctuation, symbol, or mark.

    Corrupt or blank PDFs often yield only whitespace, zero-width spaces (U+200B), or
    other format/control characters that survive :meth:`str.strip` but carry no content.
    """
    for ch in text:
        cat = unicodedata.category(ch)
        if cat[0] in ("L", "N", "P", "S", "M"):
            return True
    return False


_GLYPH_NAME_TOKEN_RE = re.compile(r"^/[A-Za-z]{0,4}\d{1,5}$")
_CID_TOKEN_RE = re.compile(r"\(cid:\d+\)")


def _letter_ratio(text: str) -> float:
    """Fraction of non-whitespace characters in *text* that are letters (any script)."""
    non_space = [c for c in text if not c.isspace()]
    if not non_space:
        return 0.0
    return sum(1 for c in non_space if c.isalpha()) / len(non_space)


def _looks_like_encoded_glyphs(
    text: str,
    *,
    min_tokens: int = 8,
    ratio_threshold: float = 0.4,
    min_len: int = 40,
    cid_char_ratio_threshold: float = 0.3,
    min_letter_ratio: float = 0.3,
) -> bool:
    """True if *text* looks like an undecoded/garbled PDF text layer rather than
    real extracted content.

    Covers three concrete failure modes seen in practice, all stemming from the
    same root cause — a font ``/Encoding`` with a ``/Differences`` array mapping
    character codes to non-standard glyph names, combined with no working
    ``/ToUnicode`` CMap, so no text-layer library can recover the real characters
    (the mapping information genuinely is not in the file; the PDF still renders
    correctly in any viewer because rendering uses the embedded font program
    directly, not this code-to-Unicode mapping):

    1. pypdf-style raw glyph-name tokens (``/143``, ``/j0``) — pypdf falls back to
       emitting the PDF name object itself when it can't resolve it.
    2. pdfminer/pdfplumber-style CID markers (``(cid:143)``) — same fallback, pdfminer's
       own notation.
    3. Byte-substitution garbage: some pdfminer fallback paths map each character
       code into some other single ASCII punctuation/digit character instead of
       raising, producing a per-glyph substitution cipher that still *looks*
       word-shaped (correct spacing) because only the space glyph happens to
       resolve correctly — no simple constant shift recovers it. Neither pattern
       above catches this, so it needs a general, format-agnostic signal: real
       prose in any language is dominated by letters even in fairly numeric
       documents (headers, labels, connecting words); text this ``[A-Za-z]``-poor
       relative to its digit/punctuation content is not real language.

    All three still count as "meaningful" for :func:`_has_meaningful_text`
    (digits/slashes/parens are letters, numbers or punctuation), so they pass
    through silently otherwise — and none of this is detectable via the presence
    of images (unlike a scanned/image-only page).
    """
    tokens = text.split()
    if len(tokens) >= min_tokens:
        glyph_like = sum(1 for t in tokens if _GLYPH_NAME_TOKEN_RE.match(t))
        if (glyph_like / len(tokens)) >= ratio_threshold:
            return True

    if len(text) >= min_len:
        cid_chars = sum(len(m) for m in _CID_TOKEN_RE.findall(text))
        if cid_chars and (cid_chars / len(text)) >= cid_char_ratio_threshold:
            return True
        if _letter_ratio(text) < min_letter_ratio:
            return True

    return False


_WORD_TOKEN_RE = re.compile(r"[^\W\d_]{2,}", re.UNICODE)

_SUBST_CIPHER_MIN_TOKENS = 80
_SUBST_CIPHER_WINDOW = 400
_SUBST_CIPHER_STEP = 200
_SUBST_CIPHER_MIDCAP_HARD = 0.05    # worst-window cipher-token rate at/above this
_SUBST_CIPHER_MIDCAP_SOFT = 0.035   # soft band needs the dict to also be poor
_SUBST_CIPHER_MIDCAP_STRONG = 0.10  # bar to spend OCR / whole-document re-OCR
_SUBST_CIPHER_DICT_HARD = 0.72      # window this Portuguese-dense is not a cipher
_SUBST_CIPHER_DICT_SOFT = 0.55

# Guard against swapping a page's text for an equivalent one on a false-positive
# detection: a recovery candidate must beat the original's PT-BR dictionary rate by
# this margin, over at least this many tokens.
_RECOVERY_MIN_CANDIDATE_TOKENS = 20
_RECOVERY_DICT_GAIN = 0.15

# poppler pdftotext per-page limits (bytes of text kept / seconds).
_POPPLER_TIMEOUT = 60
_POPPLER_MAX_PAGE_BYTES = 8_000_000


def _tokenize_words(text: str) -> List[str]:
    return _WORD_TOKEN_RE.findall(text)


def _is_cipher_token(token: str) -> bool:
    """True if *token* carries a mid-word lowercase->UPPERCASE transition that is
    not just two real words glued together.

    Real Portuguese tokens are lowercase, Capitalized (one leading cap) or ALL-CAPS.
    A per-glyph substitution cipher breaks that (``soQdaJHP``, ``LICITAgAO``,
    ``idenAﬁcadas``). Table extraction that drops the space between a unit column
    and a description (``unRolo``, ``kgChapa``) also produces the transition, so a
    token that splits at the first such point into ``<prefix><real word>`` is
    excluded — that is a layout artifact, not corruption.
    """
    if len(token) < 4:
        return False
    idx = next(
        (i for i in range(2, len(token)) if token[i].isupper() and token[i - 1].islower()),
        None,
    )
    if idx is None:
        return False
    right = token[idx:]
    if len(right) >= 3 and is_probably_portuguese(right):
        return False
    return True


def _substitution_cipher_score(
    text: str,
    *,
    window: int = _SUBST_CIPHER_WINDOW,
    step: int = _SUBST_CIPHER_STEP,
    min_tokens: int = _SUBST_CIPHER_MIN_TOKENS,
) -> float:
    """Worst sliding-window fraction of word-tokens that look cipher-shaped
    (see :func:`_is_cipher_token`). 0.0 when there is not enough text to judge."""
    try:
        tokens = _tokenize_words(text)
        n = len(tokens)
        if n < min_tokens:
            return 0.0
        flags = [1 if _is_cipher_token(t) else 0 for t in tokens]
        worst = 0.0
        for i in range(0, max(1, n - window + 1), step):
            seg = flags[i : i + window]
            if seg:
                worst = max(worst, sum(seg) / len(seg))
        return worst
    except Exception:
        return 0.0


def _looks_like_substitution_cipher(
    text: str,
    *,
    window: int = _SUBST_CIPHER_WINDOW,
    step: int = _SUBST_CIPHER_STEP,
    min_tokens: int = _SUBST_CIPHER_MIN_TOKENS,
    midcap_hard: float = _SUBST_CIPHER_MIDCAP_HARD,
    midcap_soft: float = _SUBST_CIPHER_MIDCAP_SOFT,
    dict_hard: float = _SUBST_CIPHER_DICT_HARD,
    dict_soft: float = _SUBST_CIPHER_DICT_SOFT,
) -> bool:
    """True if *text* looks like a per-glyph substitution cipher, not real prose.

    Targets a font ``/Differences`` map with no working ``/ToUnicode``: the output
    keeps word shape and letter/digit categories (so :func:`_looks_like_encoded_glyphs`
    and :func:`_has_meaningful_text` both pass), but the words are not real
    (``EÍesentadas``, ``couvocnrónto``, ``LICITAgAO``). Signal: rate of cipher-shaped
    tokens per sliding window (corruption is often per-page). A 119-doc clean-prose
    corpus from production sits at ~0; a low PT-BR dictionary hit-rate in the same
    window corroborates the softer band. Heavily tabular budgets and web-scraped
    documents can still trip this — the recovery step guards against actually
    replacing good text (see :func:`_recovery_is_improvement`).
    """
    try:
        tokens = _tokenize_words(text)
        n = len(tokens)
        if n < min_tokens:
            return False
        flags = [1 if _is_cipher_token(t) else 0 for t in tokens]
        for i in range(0, max(1, n - window + 1), step):
            seg = flags[i : i + window]
            if not seg:
                continue
            midcap = sum(seg) / len(seg)
            if midcap < midcap_soft:
                continue
            dr = dict_hit_rate(tokens[i : i + window])
            if midcap >= midcap_hard and dr < dict_hard:
                return True
            if midcap_soft <= midcap < midcap_hard and dr < dict_soft:
                return True
        return False
    except Exception:
        return False


def _text_layer_looks_broken(text: str) -> bool:
    """Glyph-name / CID / low-letter garbage OR a per-glyph substitution cipher."""
    return _looks_like_encoded_glyphs(text) or _looks_like_substitution_cipher(text)


def _recovery_is_improvement(original: str, candidate: str) -> bool:
    """Guard against cosmetic engine swaps on a false-positive detection.

    If the original is unambiguous glyph-name/CID garbage, any clean candidate wins.
    Otherwise (only the substitution-cipher heuristic flagged it) the candidate must
    be clearly more Portuguese before the page is replaced.
    """
    if _looks_like_encoded_glyphs(original):
        return True
    cand_tokens = _tokenize_words(candidate)
    if len(cand_tokens) < _RECOVERY_MIN_CANDIDATE_TOKENS:
        return False
    gain = dict_hit_rate(cand_tokens) - dict_hit_rate(_tokenize_words(original))
    return gain >= _RECOVERY_DICT_GAIN


class TextExtractor:
    """Main class for handling text extraction from various file formats."""

    def __init__(
        self,
        ocr_handler=False,
        use_aws=False,
        aws_access_key=None,
        aws_secret_key=None,
        aws_region='us-east-1',
        config: Optional[GoblinConfig] = None,
        suppress_warnings: Optional[bool] = None,
        extract_tables: bool = False,
        table_format: str = "markdown",
    ):
        """
        Initialize the text extractor.

        Args:
            ocr_handler: Enable OCR for image-based PDFs
            use_aws: Use AWS Textract for OCR
            aws_access_key: AWS access key
            aws_secret_key: AWS secret key
            aws_region: AWS region
            config: GoblinConfig object (overrides other parameters)
            suppress_warnings: If True/False, sets warning suppression for the process.
                If None (default), leaves the current setting unchanged (use
                ``goblintools.configure(suppress_warnings=True)`` at startup, or pass
                ``FileManager(suppress_warnings=True)`` before extraction).
            extract_tables: When True, detect PDF tables (pdfplumber) and embed them
                in the extracted text (currently Markdown only).
            table_format: Output format for embedded tables (``"markdown"``).
        """
        self.config = config or GoblinConfig.default()

        # Override config with explicit parameters if provided
        if any([use_aws, aws_access_key, aws_secret_key, aws_region != 'us-east-1']):
            self.config.ocr = OCRConfig(use_aws, aws_access_key, aws_secret_key, aws_region)

        if suppress_warnings is not None:
            _set_suppress_warnings(suppress_warnings)

        if table_format != "markdown":
            raise ValueError(
                f"Unsupported table_format={table_format!r}; only 'markdown' is supported"
            )
        self.extract_tables = extract_tables
        self.table_format = table_format

        if ocr_handler:
            self.ocr_handler = OCRProcessor(self.config.ocr)
        else:
            self.ocr_handler = None

        # Confidence / provenance of the most recent PDF extraction (None otherwise).
        self.last_extraction_report: Optional[ExtractionReport] = None
        # Per-file reports from the most recent ``extract_from_folder`` call,
        # keyed by the same relative path used in the ``file_path_pwd`` tag.
        self.last_extraction_reports: Dict[str, ExtractionReport] = {}

        self._parsers = None  # Lazy initialization

    @property
    def parsers(self) -> Dict[str, Callable]:
        """Lazy-loaded parsers dictionary"""
        if self._parsers is None:
            self._parsers = self._initialize_parsers()
        return self._parsers

    def _initialize_parsers(self) -> Dict[str, Callable]:
        """Initialize all available text extraction parsers."""
        return {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.txt': self._extract_txt,
            '.pptx': self._extract_pptx,
            '.html': self._extract_html,
            '.odt': self._extract_odt,
            '.rtf': self._extract_rtf,
            '.csv': self._extract_csv,
            '.xml': self._extract_xml,
            '.xlsx': self._extract_xlsx,
            '.xlsm': self._extract_xlsx,
            '.xls': self._extract_xls,
            '.ods': self._extract_ods,
            '.dbf': self._extract_dbf,
        }

    def add_parser(self, extension: str, parser_func: Callable) -> None:
        """Add or override a parser for a specific file extension."""
        self.parsers[extension.lower()] = parser_func

    def extract_from_file(self, file_path: str, display_path: Optional[str] = None) -> str:
        """
        Extract text from a single file.

        Args:
            file_path: Path to the file to extract text from
            display_path: Optional path for the file_path_pwd tag (e.g. relative path from inside archive).
                          If None, uses file_path.

        Returns:
            Extracted text as string with file_path_pwd tag at the beginning
        """
        self.last_extraction_report = None

        if not os.path.exists(file_path):
            log_warning(logger, f"File not found: {file_path}")
            return ""

        file_extension = Path(file_path).suffix.lower()
        parser = self.parsers.get(file_extension)

        if not parser:
            detected = FileValidator.detect_extension_from_magic(file_path)
            if detected:
                parser = self.parsers.get(detected)
                if parser:
                    file_extension = detected

        if not parser:
            log_warning(
                logger,
                f"No parser available for file extension: {file_extension or '(none)'}",
            )
            return ""

        try:
            extracted_text = parser(file_path)
            if not extracted_text or not _has_meaningful_text(str(extracted_text)):
                return ""
            path_for_tag = display_path if display_path is not None else file_path
            return f'file_path_pwd:"{path_for_tag}"\n{extracted_text}'

        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""

    def extract_from_folder(self, folder_path: str) -> str:
        """
        Extract text from all supported files in a folder (recursively).

        Args:
            folder_path: Path to the folder to process

        Returns:
            Combined extracted text with file_path_pwd tags for each file

        Note:
            ``last_extraction_report`` holds only the *last* PDF processed here.
            For per-file confidence use ``last_extraction_reports`` — a dict keyed
            by the same relative path as the ``file_path_pwd`` tag.
        """
        self.last_extraction_reports = {}

        if not os.path.exists(folder_path):
            log_warning(logger, f"Folder not found: {folder_path}")
            return ""

        all_texts = []
        for root, _, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                try:
                    if os.path.getsize(file_path) > self.config.max_file_size:
                        log_warning(logger, f"Skipping large file: {file_path}")
                        continue
                except OSError:
                    continue

                rel_path = os.path.relpath(file_path, folder_path)
                text = self.extract_from_file(file_path, display_path=rel_path)
                if self.last_extraction_report is not None:
                    self.last_extraction_reports[rel_path] = self.last_extraction_report
                if text:
                    all_texts.append(text)

        return '\n\n'.join(all_texts)

    def extract_tables_from_pdf(
        self, file_path: str, *, max_pages: Optional[int] = None
    ) -> List[Dict]:
        """Extract structured tables from a PDF (pdfplumber).

        Returns a list of dicts with keys ``page`` (1-based), ``index``, and ``rows``.
        """
        return extract_pdf_tables(file_path, max_pages=max_pages)

    def pdf_needs_ocr(self, file_path: str) -> bool:
        """Check if PDF needs OCR processing"""
        try:
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    try:
                        text = self._pypdf_try_extract_text(page)
                    except Exception:
                        continue
                    if text and not text.isspace():
                        return False
            return True
        except Exception as e:
            logger.error(f"Error checking PDF {file_path}: {e}")
            return True

    def _pypdf_try_extract_text(self, page) -> str:
        """Try plain then layout (newer pypdf); fall back to legacy extract_text()."""
        last_error: Optional[Exception] = None
        for mode in ("plain", "layout"):
            try:
                t = page.extract_text(extraction_mode=mode)
                if t:
                    return t
            except TypeError:
                break
            except Exception as e:
                last_error = e
                continue
        try:
            return page.extract_text() or ""
        except Exception:
            if last_error:
                raise last_error
            raise

    def _pdf_page_has_images(self, page) -> bool:
        resources = page.get("/Resources")
        if isinstance(resources, IndirectObject):
            resources = resources.get_object()
        if not resources or "/XObject" not in resources:
            return False
        xobject = resources["/XObject"]
        if isinstance(xobject, IndirectObject):
            xobject = xobject.get_object()
        return any(
            xobject[obj].get("/Subtype") == "/Image" for obj in xobject
        )

    def _pypdf_extract_pages(
        self, pdf_path: str
    ) -> Tuple[List[str], Set[int], bool]:
        """Per-page PyPDF text; failed indices; whether any page references images."""
        reader = PdfReader(pdf_path)
        texts: List[str] = []
        failed: Set[int] = set()
        has_images = False
        for i, page in enumerate(reader.pages):
            try:
                text = self._pypdf_try_extract_text(page)
                texts.append(text or "")
            except Exception as e:
                log_warning(logger, f"Error reading page {i} of {pdf_path}: {e}")
                texts.append("")
                failed.add(i)
            try:
                if not has_images and self._pdf_page_has_images(page):
                    has_images = True
            except Exception:
                pass
        return texts, failed, has_images

    def _pdfplumber_page_text(self, pdf_path: str, indices: List[int]) -> Dict[int, str]:
        """Best-effort text for specific 0-based pages via pdfplumber.

        Tried as a cheaper alternative to OCR when pypdf's font/``/Differences``
        resolution fails (see :func:`_looks_like_encoded_glyphs`): pdfplumber
        (pdfminer.six) has independent font-handling code that sometimes succeeds
        where pypdf's does not, and it's already a goblintools dependency (used for
        table extraction), so trying it first costs nothing extra to install.
        """
        out: Dict[int, str] = {}
        try:
            import pdfplumber
        except ImportError:
            return out
        try:
            with pdfplumber.open(pdf_path) as pdf:
                for idx in indices:
                    if idx < 0 or idx >= len(pdf.pages):
                        continue
                    try:
                        out[idx] = pdf.pages[idx].extract_text() or ""
                    except Exception as e:
                        log_warning(
                            logger, f"pdfplumber failed on page {idx} of {pdf_path}: {e}"
                        )
        except Exception as e:
            log_warning(logger, f"pdfplumber failed to open {pdf_path}: {e}")
        return out

    def _poppler_page_text(self, pdf_path: str, indices: List[int]) -> Dict[int, str]:
        """Best-effort text for specific 0-based pages via poppler's ``pdftotext``.

        Tried after pdfplumber and before OCR in the broken-text-layer recovery
        chain: a third independent engine (poppler, not pdfminer or pypdf) with no
        Python build dependency — just the ``pdftotext`` binary that ``pdf2image``
        already needs. One subprocess per page keeps each call's output bounded.
        """
        out: Dict[int, str] = {}
        exe = shutil.which("pdftotext")
        if not exe:
            log_warning(
                logger,
                "pdftotext (poppler-utils) not found on PATH; skipping poppler "
                "recovery step.",
            )
            return out
        for idx in indices:
            if idx < 0:
                continue
            page_no = str(idx + 1)
            try:
                proc = subprocess.run(
                    [exe, "-q", "-f", page_no, "-l", page_no, "--", pdf_path, "-"],
                    capture_output=True,
                    timeout=_POPPLER_TIMEOUT,
                    check=False,
                )
            except (subprocess.SubprocessError, OSError) as e:
                log_warning(
                    logger, f"pdftotext failed on page {idx} of {pdf_path}: {e}"
                )
                continue
            if proc.returncode != 0:
                log_warning(
                    logger,
                    f"pdftotext exited {proc.returncode} on page {idx} of {pdf_path}",
                )
                continue
            if len(proc.stdout) > _POPPLER_MAX_PAGE_BYTES:
                log_warning(
                    logger,
                    f"pdftotext output for page {idx} of {pdf_path} exceeds "
                    f"{_POPPLER_MAX_PAGE_BYTES} bytes; discarding.",
                )
                continue
            out[idx] = proc.stdout.decode("utf-8", "replace")
        return out

    def _recover_glyph_garbage_pages(
        self,
        file_path: str,
        page_texts: List[str],
        skip_indices: Set[int],
        report: Optional[ExtractionReport] = None,
    ) -> List[str]:
        """Replace pages whose text layer looks broken — glyph-name garbage OR a
        per-glyph substitution cipher.

        Chain: pdfplumber -> poppler ``pdftotext`` -> per-page OCR (if a handler is
        set). Pages in *skip_indices* (already handled as pypdf read failures) are
        left alone. Mutates and returns *page_texts*; fills *report* page statuses
        when a report is given.
        """
        broken_pages = {
            i
            for i, t in enumerate(page_texts)
            if i not in skip_indices and _text_layer_looks_broken(t)
        }
        if report is not None:
            for i in broken_pages:
                if 0 <= i < len(report.pages):
                    report.pages[i].broken_before = True
        if not broken_pages:
            return page_texts

        log_warning(
            logger,
            f"{len(broken_pages)} page(s) of {file_path} have a broken text layer "
            "(glyph-name garbage or per-glyph substitution cipher); trying "
            "pdfplumber, then poppler, then OCR.",
        )

        def _mark(idx: int, status: str, engine: str) -> None:
            if report is not None and 0 <= idx < len(report.pages):
                report.pages[idx].status = status
                report.pages[idx].engine = engine

        def _accept(idx: int, cand: str) -> bool:
            if not cand or not _has_meaningful_text(cand):
                return False
            if _text_layer_looks_broken(cand):
                return False
            return _recovery_is_improvement(page_texts[idx], cand)

        remaining = set(broken_pages)

        for engine_name, getter in (
            ("pdfplumber", self._pdfplumber_page_text),
            ("poppler", self._poppler_page_text),
        ):
            if not remaining:
                break
            try:
                got = getter(file_path, sorted(remaining))
            except Exception as e:  # never abort the caller
                log_warning(
                    logger, f"{engine_name} recovery failed for {file_path}: {e}"
                )
                got = {}
            for idx in sorted(remaining):
                cand = got.get(idx, "")
                if _accept(idx, cand):
                    page_texts[idx] = cand
                    _mark(idx, RECOVERED, engine_name)
                    remaining.discard(idx)

        # Spend OCR only on pages with strong evidence of corruption — unambiguous
        # glyph-name garbage, or a cipher rate well past the detection threshold.
        # Milder flags (heavily tabular pages, borderline corruption) are reported
        # as corrupt rather than OCR'd, since OCR of a table is usually worse.
        ocr_worthy = {
            idx
            for idx in remaining
            if _looks_like_encoded_glyphs(page_texts[idx])
            or _substitution_cipher_score(page_texts[idx]) >= _SUBST_CIPHER_MIDCAP_STRONG
        }
        if ocr_worthy and self.ocr_handler:
            logger.info(
                "OCR fallback for %d page(s) of %s (pypdf, pdfplumber and poppler "
                "all produced a broken text layer)",
                len(ocr_worthy),
                file_path,
            )
            ocr_by_page = self.ocr_handler.extract_text_from_pdf_page_indices(
                file_path, sorted(ocr_worthy)
            )
            for idx in sorted(ocr_worthy):
                cand = ocr_by_page.get(idx, "")
                if cand:
                    page_texts[idx] = cand
                    _mark(idx, RECOVERED, "ocr")
                    remaining.discard(idx)

        for idx in sorted(remaining):
            _mark(idx, CORRUPT_UNRECOVERABLE, "")
            log_warning(
                logger,
                f"page {idx} of {file_path} still has a broken text layer after "
                + (
                    "pdfplumber, poppler and OCR"
                    if self.ocr_handler
                    else "pdfplumber and poppler, and no OCR handler was provided"
                ),
            )

        return page_texts

    def _merge_page_texts(
        self, primary: List[str], secondary: List[str]
    ) -> Tuple[List[str], Set[int]]:
        """Fill empty primary slots from secondary; return merged list and still-empty indices."""
        n = max(len(primary), len(secondary))
        merged = []
        for i in range(n):
            a = primary[i] if i < len(primary) else ""
            b = secondary[i] if i < len(secondary) else ""
            merged.append(a if (a and a.strip()) else (b or ""))
        failed = {i for i, t in enumerate(merged) if not (t and t.strip())}
        return merged, failed

    def _resave_pdf(self, file_path: str) -> str:
        """Resave PDF to fix potential xref/stream issues (used as second attempt)."""
        reader = PdfReader(file_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        output_path = Path(file_path).with_suffix(".resaved.pdf")
        with open(output_path, "wb") as f:
            writer.write(f)

        return str(output_path)

    def validate_installation(self) -> Dict[str, bool]:
        """Check if all dependencies are properly installed"""
        results = {}

        # Check Tesseract
        try:
            import pytesseract
            pytesseract.get_tesseract_version()
            results['tesseract'] = True
        except Exception:
            results['tesseract'] = False

        # Check AWS credentials (use ocr_handler.use_aws in case we fell back to local)
        if self.ocr_handler and self.ocr_handler.use_aws:
            try:
                self.ocr_handler.textract_client.list_document_analysis_jobs
                results['aws_textract'] = True
            except Exception:
                results['aws_textract'] = False

        return results

    # Individual parser methods
    def _extract_pdf(self, file_path: str) -> str:
        """PyPDF text; resave retry; pdfplumber/poppler/OCR recovery for broken text
        layers; per-page OCR for gaps when a handler is set."""
        temp_file: Optional[str] = None
        page_texts: List[str] = []
        has_images = False
        extracted_text = ""
        report = ExtractionReport(path=file_path)

        def run_pypdf(path: str) -> Tuple[List[str], Set[int], bool]:
            return self._pypdf_extract_pages(path)

        try:
            try:
                page_texts, failed, has_images = run_pypdf(file_path)
            except Exception as e:
                log_warning(
                    logger,
                    f"PyPDF read failed on original {file_path}: {e}; trying resaved copy.",
                )
                page_texts, failed, has_images = [], set(), False
                temp_file = self._resave_pdf(file_path)
                page_texts, failed, has_images = run_pypdf(temp_file)
            else:
                if failed:
                    try:
                        temp_file = self._resave_pdf(file_path)
                        alt_texts, _alt_failed, alt_img = run_pypdf(temp_file)
                        page_texts, still_failed = self._merge_page_texts(
                            page_texts, alt_texts
                        )
                        has_images = has_images or alt_img
                        failed = still_failed
                    except Exception as e:
                        log_warning(
                            logger,
                            f"PyPDF resave retry failed for {file_path}: {e}",
                        )

            report.pages = [
                PageExtraction(index=i, engine="pypdf")
                for i in range(len(page_texts))
            ]

            if failed and self.ocr_handler:
                logger.info(
                    "OCR fallback for %d page(s) of %s (PyPDF could not decode them)",
                    len(failed),
                    file_path,
                )
                ocr_by_page = self.ocr_handler.extract_text_from_pdf_page_indices(
                    file_path, sorted(failed)
                )
                for idx in failed:
                    if idx in ocr_by_page and ocr_by_page[idx]:
                        page_texts[idx] = ocr_by_page[idx]
                        if 0 <= idx < len(report.pages):
                            report.pages[idx].status = RECOVERED
                            report.pages[idx].engine = "ocr"

            page_texts = self._recover_glyph_garbage_pages(
                file_path, page_texts, failed, report
            )

            if self.extract_tables and page_texts:
                page_texts = self._merge_tables_into_page_texts(file_path, page_texts)

            extracted_text = "\n".join(page_texts)

        except Exception as e:
            logger.error(f"Failed to open PDF {file_path}: {e}")
            self.last_extraction_report = None
            return ""
        finally:
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except OSError:
                    pass

        for pe, txt in zip(report.pages, page_texts):
            if pe.status == CLEAN and not (txt and txt.strip()):
                pe.status = CORRUPT_UNRECOVERABLE
                pe.engine = ""

        # Whole-document safety net: catches pages that individually stayed under the
        # per-page thresholds but still add up to unreadable output, the pre-existing
        # blank/corrupt-page case, and a whole-document substitution cipher. Not gated
        # on has_images: OCR rasterizes regardless, and a broken font /Differences
        # encoding has nothing to do with embedded images.
        if (
            not _has_meaningful_text(extracted_text)
            or _looks_like_encoded_glyphs(extracted_text)
            or _substitution_cipher_score(extracted_text) >= _SUBST_CIPHER_MIDCAP_STRONG
        ):
            if self.ocr_handler:
                logger.info(f"OCR required for file: {file_path}")
                ocr_text = self.ocr_handler.extract_text_from_pdf(file_path)
                report.used_ocr = True
                for pe in report.pages:
                    if pe.status in (CLEAN, CORRUPT_UNRECOVERABLE):
                        pe.status = RECOVERED
                        pe.engine = "ocr"
                report.recompute_overall()
                self.last_extraction_report = report
                self._log_report_warning(report)
                return ocr_text
            log_warning(
                logger,
                f"The file {file_path} requires OCR but no handler was provided.",
            )
            for pe in report.pages:
                if pe.status == CLEAN:
                    pe.status = CORRUPT_UNRECOVERABLE
                    pe.engine = ""

        report.recompute_overall()
        self.last_extraction_report = report
        self._log_report_warning(report)
        return extracted_text

    def _log_report_warning(self, report: ExtractionReport) -> None:
        """Emit one structured warning when the last extraction is not fully clean."""
        if report.is_clean:
            return
        n_recovered = sum(1 for p in report.pages if p.status == RECOVERED)
        n_corrupt = sum(1 for p in report.pages if p.status == CORRUPT_UNRECOVERABLE)
        by_engine: Dict[str, int] = {}
        for p in report.pages:
            if p.status == RECOVERED and p.engine:
                by_engine[p.engine] = by_engine.get(p.engine, 0) + 1
        detail = ", ".join(f"{k}:{v}" for k, v in sorted(by_engine.items())) or "none"
        log_warning(
            logger,
            f"{report.path}: text-layer confidence '{report.overall_status}' "
            f"({len(report.pages)} pages, {n_recovered} recovered [{detail}], "
            f"{n_corrupt} still corrupt)",
        )

    def _merge_tables_into_page_texts(
        self, file_path: str, page_texts: List[str]
    ) -> List[str]:
        """Append Markdown tables after each page's pypdf text."""
        try:
            tables = extract_pdf_tables(file_path, max_pages=len(page_texts))
        except ImportError as e:
            log_warning(logger, str(e))
            return page_texts
        except Exception as e:
            log_warning(logger, f"Table extraction skipped for {file_path}: {e}")
            return page_texts

        if not tables:
            return page_texts

        by_page = tables_by_page(tables)
        merged: List[str] = []
        for i, text in enumerate(page_texts):
            page_num = i + 1
            page_tables = by_page.get(page_num, [])
            if not page_tables:
                merged.append(text)
                continue
            block = format_tables_for_page(
                page_tables, page=page_num, table_format=self.table_format
            )
            if block:
                base = text.rstrip()
                merged.append(f"{base}\n\n{block}" if base else block)
            else:
                merged.append(text)
        return merged

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        try:
            doc = docx.Document(file_path)
            return ' '.join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")
            return ""

    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error processing TXT file {file_path}: {e}")
                return ""

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        try:
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            return ""

    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML files."""
        encodings = ['utf-8', 'latin-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    soup = BeautifulSoup(file.read(), 'html.parser')
                    return soup.get_text(separator=' ', strip=True)
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error processing HTML file {file_path}: {e}")
                break
        return ""

    def _extract_odt(self, file_path: str) -> str:
        """Extract text from OpenDocument Text files."""
        try:
            doc = load(file_path)
            return ' '.join(
                teletype.extractText(element)
                for element in doc.getElementsByType(text.P)
            )
        except Exception as e:
            logger.error(f"Error processing ODT file {file_path}: {e}")
            return ""

    def _extract_rtf(self, file_path: str) -> str:
        """Extract text from RTF files."""
        try:
            with open(file_path, 'r') as file:
                return rtf_to_text(file.read(), errors='ignore')
        except Exception as e:
            logger.error(f"Error processing RTF file {file_path}: {e}")
            return ""

    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return ' '.join(
                    ' '.join(row)
                    for row in csv.reader(file)
                    if any(field.strip() for field in row)
                )
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {e}")
            return ""

    def _extract_xml(self, file_path: str) -> str:
        """Extract text from XML files."""
        try:
            tree = ET.parse(file_path)
            return ' '.join(
                elem.text.strip()
                for elem in tree.iter()
                if elem.text and elem.text.strip()
            )
        except Exception as e:
            logger.error(f"Error processing XML file {file_path}: {e}")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract evaluated text content from Excel files."""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            texts.append(str(cell))
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing XLSX file {file_path}: {e}")
            return ""

    def _extract_xls(self, file_path: str) -> str:
        """Extract text from legacy Excel files."""
        try:
            book = xlrd.open_workbook(file_path, formatting_info=False)
            texts = []
            for sheet in book.sheets():
                for row_idx in range(sheet.nrows):
                    for cell in sheet.row(row_idx):
                        value = cell.value
                        if value and not str(value).startswith('='):
                            texts.append(str(value))
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing XLS file {file_path}: {e}")
            return ""

    def _extract_ods(self, file_path: str) -> str:
        """Extract text from OpenDocument Spreadsheets."""
        try:
            doc = load(file_path)
            return '\n'.join(
                "".join(
                    child.data
                    for child in p.childNodes
                    if child.nodeType == child.TEXT_NODE
                )
                for p in doc.getElementsByType(P)
            )
        except Exception as e:
            logger.error(f"Error processing ODS file {file_path}: {e}")
            return ""

    def _extract_dbf(self, file_path: str) -> str:
        """Extract text from DBF database files."""
        try:
            dbf = DBF(file_path, load=True)
            return ' '.join(
                f"{key}: {value}"
                for record in dbf
                for key, value in record.items()
            )
        except Exception as e:
            logger.error(f"Error processing DBF file {file_path}: {e}")
            return ""
